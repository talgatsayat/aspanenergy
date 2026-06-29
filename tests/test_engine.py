"""Calibration tests — verify the engine reproduces the real NESKAO offer.

Reference figures from the Aspan NESKAO 600 kWc offer & technical report:
  - 600 kWc recommended, ~1000 modules
  - 870,000 kWh/yr production (1450 kWh/kWc)
  - Solar tariff 67.10 FCFA/kWh (83.87 x 0.80)
  - Scenario 1 unit saving 16.77 FCFA/kWh -> ~14.6 M/yr -> ~219 M over 15 yr
Run with:  python -m pytest -q   (or  python tests/test_engine.py)
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from aspan.engine import (  # noqa: E402
    build_proposal_data,
    compute_scenarios,
    size_system,
    summarize_bills,
)
from aspan.geo import panels_from_area, parse_dms  # noqa: E402
from aspan.schema import ClientInfo, MonthlyBill  # noqa: E402


def _neskao_bills():
    # Representative NESKAO monthly consumption ~256,800 kWh, day tariff 83.87.
    # Daytime share derived from CIE bill (Jour+Pointe)/Total ~0.68.
    return [
        MonthlyBill(period=f"2025.{m:02d}", total_kwh=256_816, day_tariff=83.87)
        for m in range(1, 13)
    ]


def _neskao_client():
    return ClientInfo(
        name="NESKAO",
        industry="Food & beverage manufacturing",
        country="Cote d'Ivoire",
        has_diesel=True,
        grid_capacity_kva=1250,
        roof_area_m2=5149,  # sum of T9 + cold room + cleaning + roasting roofs
    )


def approx(a, b, tol=0.03):
    return abs(a - b) <= tol * abs(b)


def test_panels_from_area():
    est = panels_from_area(5149)
    # ~1000 modules, ~600 kWc on NESKAO roofs
    assert 950 <= est.n_panels <= 1010, est.n_panels
    assert approx(est.capacity_kwc, 600), est.capacity_kwc


def test_parse_dms():
    assert approx(parse_dms("5°15'53.8\"N"), 5.2649, 0.01)
    assert parse_dms("4°00'12.4\"W") < 0  # West is negative


def test_sizing_recommends_600():
    summary = summarize_bills(_neskao_bills())
    sizing = size_system(_neskao_client(), summary)
    assert approx(sizing.recommended_kwc, 600), sizing.recommended_kwc
    assert sizing.binding_constraint == "roof"
    assert approx(sizing.annual_production_kwh, 870_000), sizing.annual_production_kwh


def test_scenario1_savings():
    summary = summarize_bills(_neskao_bills())
    sizing = size_system(_neskao_client(), summary)
    s1, s2 = compute_scenarios(sizing, summary)

    assert approx(s1.solar_tariff_per_kwh, 67.10)
    assert approx(s1.saving_per_kwh, 16.77)
    assert approx(s1.annual_saving, 14_600_000, 0.03), s1.annual_saving
    assert approx(s1.cumulative_15y_saving, 219_000_000, 0.03), s1.cumulative_15y_saving


def test_scenario2_diesel():
    summary = summarize_bills(_neskao_bills())
    sizing = size_system(_neskao_client(), summary)
    _, s2 = compute_scenarios(sizing, summary)
    # baseline = 0.5*83.87 + 0.5*210 = 146.935 ; saving = 79.835 FCFA/kWh
    assert approx(s2.baseline_cost_per_kwh, 146.94, 0.01)
    assert approx(s2.saving_per_kwh, 79.84, 0.01)
    assert s2.annual_saving > s1_annual_for_compare()


def s1_annual_for_compare():
    summary = summarize_bills(_neskao_bills())
    sizing = size_system(_neskao_client(), summary)
    s1, _ = compute_scenarios(sizing, summary)
    return s1.annual_saving


def test_end_to_end():
    data = build_proposal_data(_neskao_client(), _neskao_bills(), reference="OF-NESKAO-TEST")
    assert data.sizing.recommended_kwc > 0
    assert len(data.scenarios) == 2
    assert data.solar_tariff < data.cie_day_tariff


def test_financials():
    data = build_proposal_data(_neskao_client(), _neskao_bills())
    f = data.financial
    # NPV must be positive and below the undiscounted 15-yr total
    assert f.npv_savings > 0
    assert f.npv_savings < data.scenarios[0].cumulative_15y_saving
    # escalation makes the escalated NPV larger than the flat NPV
    assert f.npv_savings_escalated > f.npv_savings
    # coverage between 0 and 100%
    assert 0 < f.solar_coverage_pct <= 100
    assert len(f.cashflow) == data.contract_years
    assert f.co2_tonnes_per_year > 0


def test_pvgis_yield_injection():
    # A higher site-specific yield must increase production proportionally.
    base = build_proposal_data(_neskao_client(), _neskao_bills())
    hi = build_proposal_data(_neskao_client(), _neskao_bills(),
                             specific_yield=1600, yield_source="PVGIS @ test")
    assert hi.sizing.yield_source == "PVGIS @ test"
    assert hi.sizing.specific_yield == 1600
    # 600 kWc * 1600 = 960,000 kWh
    assert approx(hi.sizing.annual_production_kwh, 600 * 1600, 0.001)
    assert hi.sizing.annual_production_kwh > base.sizing.annual_production_kwh


def test_site_report_adds_feasibility_slide(tmp_path=None):
    import os
    import tempfile

    from aspan.pptx_builder import build_pptx
    from aspan.schema import BuildingRoof, SiteReport
    from pptx import Presentation

    rep = SiteReport(
        latitude=5.2649, longitude=-4.0034, transformer_kva=1250, diesel_kw=1100,
        recommended_kwc=600, n_modules=1000,
        buildings=[BuildingRoof("T9", 2152, 320), BuildingRoof("Cold room", 768, 130)],
        feasibility_verdict="Project technically feasible.", engineer_name="Eng. X")

    base = build_proposal_data(_neskao_client(), _neskao_bills())
    withrep = build_proposal_data(_neskao_client(), _neskao_bills(), site_report=rep)
    assert withrep.site_report is not None

    out_base = os.path.join(tempfile.mkdtemp(), "base.pptx")
    out_rep = os.path.join(tempfile.mkdtemp(), "rep.pptx")
    build_pptx(base, out_base)
    build_pptx(withrep, out_rep)
    n_base = len(Presentation(out_base).slides)
    n_rep = len(Presentation(out_rep).slides)
    assert n_rep == n_base + 1  # exactly one extra (feasibility) slide


def test_cover_email():
    from aspan.email_gen import build_cover_email

    data = build_proposal_data(_neskao_client(), _neskao_bills())
    email = build_cover_email(data, "fr", use_ai=False)
    assert email["mode"] == "template"
    assert "NESKAO" in email["subject"]
    assert "600" in email["body"]  # mentions the recommended kWc
    assert email["subject"]


if __name__ == "__main__":
    import traceback

    tests = [v for k, v in sorted(globals().items()) if k.startswith("test_")]
    passed = 0
    for t in tests:
        try:
            t()
            print(f"PASS  {t.__name__}")
            passed += 1
        except Exception:
            print(f"FAIL  {t.__name__}")
            traceback.print_exc()
    print(f"\n{passed}/{len(tests)} tests passed")
    sys.exit(0 if passed == len(tests) else 1)
