"""Fill the OFFICIAL Aspan deck (their .pptx) with a client's data.

Instead of redrawing slides, this opens Aspan's real presentation
(``assets/template.pptx``) and replaces only the dynamic values — client name,
capacities, tariffs, savings, reference, date — leaving every company detail,
the design, the icons and the layout exactly as Aspan made them.

The replacements are value-based: the template's NESKAO figures are swapped for
the figures computed for the current client. For the NESKAO sample the output is
therefore identical to Aspan's original deck (a good correctness check).
"""
from __future__ import annotations

import datetime
import re
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation

from .config import assumptions
from .format_utils import fmt_int, t
from .schema import ProposalData

TEMPLATE = Path(__file__).resolve().parent / "assets" / "template.pptx"
EN_DASH = "–"
APPROX = "≈"
BULLET = "•"

_MONTHS = {
    "fr": ["Janvier", "Février", "Mars", "Avril", "Mai", "Juin", "Juillet",
           "Août", "Septembre", "Octobre", "Novembre", "Décembre"],
    "en": ["January", "February", "March", "April", "May", "June", "July",
           "August", "September", "October", "November", "December"],
}


def _fr(x: float, dec: int) -> str:
    return f"{x:.{dec}f}".replace(".", ",")


def _money_short(value: float) -> str:
    """Aspan-style compact money WITHOUT currency: 14.6M -> '14,6 M', 219M -> '219 M'."""
    m = value / 1_000_000
    if abs(m) >= 1000:
        return f"{_fr(m / 1000, 2)} Md"
    if abs(m) < 20:
        return f"{_fr(m, 1)} M"
    return f"{m:.0f} M"


def _mix_amount(solar_kwh: float, grid_share: float, day: float, diesel: float,
                solar: float, years: int) -> str:
    """One illustrative grid/diesel mix line, in Aspan's exact format."""
    baseline = grid_share * day + (1 - grid_share) * diesel
    annual = solar_kwh * (baseline - solar)
    return (f"{APPROX} {_money_short(annual)} / an   {BULLET}   "
            f"{APPROX} {_money_short(annual * years)} sur {years} ans")


def _replacements(data: ProposalData, city: str) -> List[Tuple[str, str]]:
    a = assumptions()
    d = data
    s1 = d.scenarios[0]
    yld = d.sizing.specific_yield
    prod = d.sizing.annual_production_kwh
    kwc = d.sizing.recommended_kwc
    day = d.cie_day_tariff
    solar = d.solar_tariff
    disc = d.discount_pct
    years = d.contract_years
    dprice = a["diesel"]["fuel_price_fcfa_per_litre"]
    lpkwh = a["diesel"]["litres_per_kwh"]
    dcost = dprice * lpkwh
    dunit = dcost - solar
    dpct = (1 - solar / dcost) * 100 if dcost else 0
    solar_used = s1.solar_kwh_used

    unit1 = s1.saving_per_kwh
    ann1 = s1.annual_saving
    cum1 = s1.cumulative_15y_saving

    # Ordered longest/most-specific first to avoid partial clashes.
    return [
        ("OF-NESKAO-2026-001", d.reference),
        ("1 450 kWh/kWc/an à Abidjan", f"{fmt_int(yld)} kWh/kWc/an à {city}"),
        ("870 000 kWh/an", f"{fmt_int(prod)} kWh/an"),
        ("600 kWc", f"{fmt_int(kwc)} kWc"),
        ("700 FCFA/litre", f"{dprice:.0f} FCFA/litre"),
        ("0,30 L/kWh", f"{_fr(lpkwh, 2)} L/kWh"),
        ("210 FCFA/kWh", f"{dcost:.0f} FCFA/kWh"),
        ("142,90 FCFA", f"{_fr(dunit, 2)} FCFA"),
        (f"{APPROX} 31 M / an   {BULLET}   {APPROX} 466 M sur 15 ans",
         _mix_amount(solar_used, 0.85, day, dcost, solar, years)),
        (f"{APPROX} 53 M / an   {BULLET}   {APPROX} 795 M sur 15 ans",
         _mix_amount(solar_used, 0.65, day, dcost, solar, years)),
        (f"{EN_DASH} 68 %", f"{EN_DASH} {dpct:.0f} %"),
        ("16,77", _fr(unit1, 2)),
        ("14,6 M", _money_short(ann1)),
        ("219 M", _money_short(cum1)),
        ("83,87", _fr(day, 2)),
        ("67,10", _fr(solar, 2)),
        ("0,80", _fr(1 - disc / 100, 2)),
        ("PPA 15.", f"PPA {years} ans."),
        ("Juin 2026", _date_label(d.client.language)),
    ]


def _date_label(lang: str) -> str:
    now = datetime.date.today()
    return f"{_MONTHS[lang][now.month - 1]} {now.year}"


def _apply_runs(text_frame, repl, disc: int, years: int, client: str):
    for para in text_frame.paragraphs:
        for run in para.runs:
            txt = run.text
            for old, new in repl:
                if old in txt:
                    txt = txt.replace(old, new)
            # discount "20 %" (any spacing) and contract "15 ans"
            txt = re.sub(r"20(\s*)%", f"{disc}\\1%", txt)
            txt = re.sub(r"15(\s*)ans", f"{years}\\1ans", txt)
            # client name last (after the reference token was handled)
            if "NESKAO" in txt:
                txt = txt.replace("NESKAO", client)
            if txt != run.text:
                run.text = txt


def build_from_template(data: ProposalData, output_path: str,
                        city: Optional[str] = None) -> str:
    """Open Aspan's deck and fill it with this client's numbers. Returns path."""
    if not TEMPLATE.exists():
        raise FileNotFoundError(
            f"Template not found at {TEMPLATE}. Place Aspan's .pptx there.")
    prs = Presentation(str(TEMPLATE))
    repl = _replacements(data, city or "Abidjan")
    disc = data.discount_pct
    years = data.contract_years
    client = data.client.name or "Client"

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _apply_runs(shape.text_frame, repl, disc, years, client)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _apply_runs(cell.text_frame, repl, disc, years, client)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
