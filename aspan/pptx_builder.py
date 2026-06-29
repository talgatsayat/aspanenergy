"""Render a branded, bilingual Aspan Energy solar PPA proposal (.pptx).

Visual identity matched to the official Aspan Energy CI offer deck:
  - "sandwich" layout: deep forest-green cover & closing, white content slides
  - Cambria serif headings, Calibri body
  - light mint cards, medium-green stat numbers, one dark-green card with a gold
    number for emphasis
  - sun logo + decorative green circles on the dark slides
Output is a fully editable PowerPoint (native text + chart + tables).
"""
from __future__ import annotations

import datetime
import math
import os
from pathlib import Path
from typing import Optional

ICON_DIR = Path(__file__).resolve().parent / "assets" / "icons"

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .config import assumptions, branding
from .format_utils import fmt_int, fmt_kwh, fmt_money, t
from .schema import ProposalData

_MONTHS = {
    "fr": ["Janvier", "Février", "Mars", "Avril", "Mai", "Juin", "Juillet",
           "Août", "Septembre", "Octobre", "Novembre", "Décembre"],
    "en": ["January", "February", "March", "April", "May", "June", "July",
           "August", "September", "October", "November", "December"],
}


def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


class DeckBuilder:
    def __init__(self, data: ProposalData):
        self.data = data
        self.lang = data.client.language
        self.b = branding()
        c = self.b["colors"]
        # palette
        self.GREEN = c["green"]
        self.GREEN_DARK = c.get("green_dark", "1A2B22")
        self.DECOR = c.get("decor", "1C4F31")
        self.WHITE = c["text_light"]
        self.MINT = c["mint"]
        self.BORDER = c.get("border", "D9EBDF")
        self.NUM = c["number"]
        self.GOLD = c["gold"]
        self.MUTED = c["muted"]
        self.WARN = c.get("warning", "B85042")
        self.HFONT = self.b["fonts"]["heading"]
        self.BFONT = self.b["fonts"]["body"]

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.W, self.H = 13.333, 7.5
        self.blank = self.prs.slide_layouts[6]
        self._pageno = 1

        now = datetime.date.today()
        self.date_label = f"{_MONTHS[self.lang][now.month - 1]} {now.year}"

    # ---------------- low-level helpers ----------------
    def _slide(self, dark: bool = False, decor: Optional[str] = None):
        s = self.prs.slides.add_slide(self.blank)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(self.GREEN if dark else self.WHITE)
        if decor:
            self._decor(s, decor)
        return s

    def _decor(self, s, kind: str):
        circles = {
            "cover": [(8.6, -2.1, 7.0), (-1.9, 4.7, 3.7)],
            "closing": [(8.7, 3.2, 6.6)],
        }.get(kind, [])
        for x, y, d in circles:
            shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                     Inches(d), Inches(d))
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(self.DECOR)
            shp.line.fill.background()
            shp.shadow.inherit = False

    def _sun(self, s, cx: float, cy: float, body_d: float = 0.5):
        r = body_d / 2
        for k in range(8):
            ang = math.radians(k * 45)
            dist = r + 0.13
            rx, ry = cx + dist * math.cos(ang), cy + dist * math.sin(ang)
            ray = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(rx - 0.055),
                                     Inches(ry - 0.055), Inches(0.11), Inches(0.11))
            ray.fill.solid()
            ray.fill.fore_color.rgb = _rgb(self.GOLD)
            ray.line.fill.background()
            ray.shadow.inherit = False
        body = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r),
                                  Inches(body_d), Inches(body_d))
        body.fill.solid()
        body.fill.fore_color.rgb = _rgb(self.GOLD)
        body.line.fill.background()
        body.shadow.inherit = False

    def _text(self, s, x, y, w, h, text, size=16, color=None, bold=False,
              align=PP_ALIGN.LEFT, font=None, anchor=MSO_ANCHOR.TOP,
              italic=False, line_spacing=1.0, wrap=True):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            r = p.add_run()
            r.text = line
            f = r.font
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.name = font or self.BFONT
            f.color.rgb = _rgb(color or self.MUTED)
        return tb

    def _card(self, s, x, y, w, h, fill=None, radius=0.06):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill or self.MINT)
        shp.line.fill.background()
        shp.shadow.inherit = False
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
        return shp

    def _disc(self, s, x, y, d, fill, label="", txt_color=None, size=14, font=None):
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                 Inches(d), Inches(d))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
        shp.line.fill.background()
        shp.shadow.inherit = False
        if label:
            tf = shp.text_frame
            for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                setattr(tf, m, 0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = label
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.name = font or self.BFONT
            r.font.color.rgb = _rgb(txt_color or self.GREEN)
        return shp

    def _picture(self, s, cx, cy, size, img_name):
        """Place a square icon centered at (cx, cy). Returns True if drawn."""
        path = ICON_DIR / img_name
        if not path.exists():
            return False
        s.shapes.add_picture(str(path), Inches(cx - size / 2), Inches(cy - size / 2),
                             width=Inches(size), height=Inches(size))
        return True

    def _icon_badge(self, s, x, y, d, img_name, fill=None, fallback=""):
        """White (or given) circle with the icon centered on it."""
        center = (x + d / 2, y + d / 2)
        self._disc(s, x, y, d, fill or self.WHITE)
        if not self._picture(s, center[0], center[1], d * 0.56, img_name):
            # graceful fallback: short label in the circle
            self._text(s, x, y, d, d, fallback, size=12, bold=True,
                       color=self.GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    def _stat_tile(self, s, x, y, w, h, value, label, emphasis=False):
        bg = self.GREEN if emphasis else self.MINT
        self._card(s, x, y, w, h, fill=bg)
        num_color = self.GOLD if emphasis else self.NUM
        lbl_color = self.WHITE if emphasis else self.MUTED
        chars = len(str(value))
        if chars <= 7:
            vsize = 32
        elif chars <= 10:
            vsize = 26
        elif chars <= 13:
            vsize = 21
        else:
            vsize = 18
        self._text(s, x + 0.25, y + 0.18, w - 0.5, h * 0.5, value, size=vsize,
                   color=num_color, bold=True, font=self.HFONT,
                   anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        self._text(s, x + 0.27, y + h * 0.64, w - 0.5, h * 0.32, label, size=12.5,
                   color=lbl_color, anchor=MSO_ANCHOR.TOP)

    def _heading(self, s, title, dark=False, kicker=None):
        self._text(s, 0.7, 0.5, 11.9, 1.0, title, size=33, bold=True,
                   font=self.HFONT, color=self.WHITE if dark else self.GREEN)

    def _footer(self, s, dark=False):
        self._pageno += 1
        d = self.data
        label = t(self.lang,
                  f"{self.b['company']['name']}  -  Offre {d.client.name}  -  {self.date_label}",
                  f"{self.b['company']['name']}  -  {d.client.name} Proposal  -  {self.date_label}")
        col = "C9D9CE" if dark else self.MUTED
        self._text(s, 0.7, self.H - 0.5, 10.5, 0.35, label, size=9, color=col)
        self._text(s, self.W - 1.2, self.H - 0.5, 0.5, 0.35, str(self._pageno),
                   size=9, color=col, align=PP_ALIGN.RIGHT)

    def _set_cell(self, cell, text, size=11, color=None, bold=False,
                  fill=None, align=PP_ALIGN.LEFT):
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.09)
        cell.margin_right = Inches(0.09)
        cell.margin_top = 0
        cell.margin_bottom = 0
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = self.BFONT
        r.font.color.rgb = _rgb(color or self.GREEN)

    # ---------------- slides ----------------
    def cover(self):
        s = self._slide(dark=True, decor="cover")
        d = self.data
        c = self.b["company"]
        if not self._picture(s, 1.15, 1.12, 0.74, "image1.png"):
            self._sun(s, 1.15, 1.12, 0.52)
        self._text(s, 1.75, 0.88, 8, 0.5, c["name"], size=20, bold=True,
                   color=self.WHITE, font=self.HFONT)
        title = t(self.lang, "OFFRE FINANCIÈRE", "FINANCIAL PROPOSAL")
        self._text(s, 0.9, 2.35, 11.5, 1.2, title, size=56, bold=True,
                   font=self.HFONT, color=self.WHITE)
        sub = t(self.lang,
                f"Centrale solaire photovoltaïque de {fmt_int(d.sizing.recommended_kwc)} kWc "
                f"en autoconsommation",
                f"{fmt_int(d.sizing.recommended_kwc)} kWc solar PV plant for self-consumption")
        self._text(s, 0.92, 3.6, 11.3, 0.7, sub, size=20, color="D9E5DD")
        # client line (NESKAO in gold)
        self._text(s, 0.92, 4.45, 2.0, 0.4,
                   t(self.lang, "Client :", "Client:"), size=16, color="D9E5DD")
        self._text(s, 1.85, 4.45, 3.5, 0.4, d.client.name, size=16, bold=True,
                   color=self.GOLD)
        self._text(s, 4.0, 4.45, 8.0, 0.4,
                   t(self.lang, f"Contrat PPA - {d.contract_years} ans     Ref. {d.reference}",
                     f"PPA Contract - {d.contract_years} years     Ref. {d.reference}"),
                   size=16, color="D9E5DD")
        self._text(s, 0.92, 5.9, 11.3, 0.5,
                   t(self.lang, c["tagline_fr"], c["tagline_en"]),
                   size=16, color=self.GOLD, italic=True, bold=True)

    def executive_summary(self):
        s = self._slide()
        d = self.data
        s1 = d.scenarios[0]
        self._heading(s, t(self.lang, "Synthèse", "Executive summary"))
        intro = t(self.lang,
                  f"ASPAN ENERGY propose à {d.client.name} d'installer et d'exploiter une centrale "
                  f"solaire de {fmt_int(d.sizing.recommended_kwc)} kWc en autoconsommation, sans "
                  f"aucun investissement. Vous ne payez que l'énergie solaire consommée, à un tarif "
                  f"garanti {d.discount_pct}% inférieur au tarif de jour CIE.",
                  f"ASPAN ENERGY proposes to install and operate a {fmt_int(d.sizing.recommended_kwc)} kWc "
                  f"self-consumption solar plant for {d.client.name}, with zero investment. You only pay "
                  f"for the solar energy consumed, at a guaranteed tariff {d.discount_pct}% below the CIE "
                  f"day tariff.")
        self._text(s, 0.7, 1.7, 11.9, 1.2, intro, size=15.5, color=self.MUTED,
                   line_spacing=1.2)
        tiles = [
            (f"{fmt_int(d.sizing.recommended_kwc)} kWc",
             t(self.lang, "Puissance recommandée", "Recommended capacity"), False),
            (fmt_kwh(d.sizing.annual_production_kwh),
             t(self.lang, "Production solaire / an", "Solar production / yr"), False),
            (fmt_money(s1.annual_saving, self.lang),
             t(self.lang, "Économie annuelle (Sc.1)", "Annual saving (Sc.1)"), False),
            (fmt_money(s1.cumulative_15y_saving, self.lang),
             t(self.lang, f"Économie cumulée {d.contract_years} ans", f"{d.contract_years}-yr cumulative"), True),
        ]
        x, w, gap = 0.7, 2.85, 0.18
        for val, label, emph in tiles:
            self._stat_tile(s, x, 3.35, w, 1.75, val, label, emphasis=emph)
            x += w + gap
        self._text(s, 0.7, 5.55, 11.9, 1.0,
                   t(self.lang,
                     f"Tarif solaire garanti : {d.solar_tariff:.2f} FCFA/kWh, soit -{d.discount_pct}% "
                     f"sur le tarif de jour CIE ({d.cie_day_tariff:.2f} FCFA/kWh).",
                     f"Guaranteed solar tariff: {d.solar_tariff:.2f} FCFA/kWh, i.e. -{d.discount_pct}% "
                     f"on the CIE day tariff ({d.cie_day_tariff:.2f} FCFA/kWh)."),
                   size=16, color=self.NUM, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        self._footer(s)

    def current_situation(self):
        s = self._slide()
        d = self.data
        bs = d.bill_summary
        self._heading(s, t(self.lang, "Situation énergétique actuelle",
                           "Current energy situation"))
        tiles = [
            (fmt_kwh(bs.avg_monthly_kwh),
             t(self.lang, "Consommation moyenne / mois", "Avg consumption / month")),
            (f"{bs.avg_day_tariff:.2f}",
             t(self.lang, "Tarif de jour CIE (FCFA/kWh)", "CIE day tariff (FCFA/kWh)")),
            (fmt_money(bs.annual_cost_estimate, self.lang),
             t(self.lang, "Facture annuelle estimée", "Estimated annual bill")),
            (f"{bs.months_count}",
             t(self.lang, "Factures analysées", "Bills analysed")),
        ]
        x, w, gap = 0.7, 2.85, 0.18
        for val, label in tiles:
            self._stat_tile(s, x, 1.75, w, 1.6, val, label)
            x += w + gap
        day_share = bs.annual_day_kwh / bs.annual_kwh if bs.annual_kwh else 0.68
        self._text(s, 0.7, 3.75, 11.9, 0.4,
                   t(self.lang, "Répartition jour / nuit de la consommation",
                     "Day / night split of consumption"),
                   size=14, bold=True, color=self.GREEN)
        bx, by, bw, bh = 0.7, 4.25, 11.9, 0.55
        dw = bw * day_share
        self._card(s, bx, by, dw, bh, fill=self.GOLD, radius=0.5)
        self._text(s, bx + 0.2, by, max(dw - 0.4, 1.4), bh,
                   t(self.lang, f"Jour {day_share*100:.0f}%", f"Day {day_share*100:.0f}%"),
                   size=12, bold=True, color=self.GREEN, anchor=MSO_ANCHOR.MIDDLE)
        self._card(s, bx + dw + 0.06, by, bw - dw - 0.06, bh, fill=self.BORDER, radius=0.5)
        self._text(s, bx + dw + 0.25, by, bw - dw - 0.5, bh,
                   t(self.lang, f"Nuit {(1-day_share)*100:.0f}%", f"Night {(1-day_share)*100:.0f}%"),
                   size=12, bold=True, color=self.GREEN, anchor=MSO_ANCHOR.MIDDLE)
        note = t(self.lang,
                 "Le solaire couvre la consommation de journée, là où se concentre l'essentiel "
                 "de l'activite industrielle. " +
                 ("Un groupe électrogène assure le secours en cas de coupure."
                  if d.client.has_diesel else ""),
                 "Solar covers daytime consumption, where most industrial activity occurs. " +
                 ("A diesel generator provides backup during grid outages."
                  if d.client.has_diesel else ""))
        self._text(s, 0.7, 5.25, 11.9, 1.0, note, size=14, color=self.MUTED,
                   line_spacing=1.2)
        self._footer(s)

    def solar_solution(self):
        s = self._slide()
        d = self.data
        sz = d.sizing
        self._heading(s, t(self.lang, "Solution solaire proposée",
                           "Proposed solar solution"))
        tiles = [
            (f"{fmt_int(sz.recommended_kwc)} kWc",
             t(self.lang, "Puissance installée", "Installed capacity")),
            (fmt_int(sz.n_panels),
             t(self.lang, "Modules photovoltaïques", "PV modules")),
            (fmt_kwh(sz.annual_production_kwh),
             t(self.lang, "Production estimée / an", "Estimated production / yr")),
        ]
        x, w, gap = 0.7, 3.83, 0.2
        for val, label in tiles:
            self._stat_tile(s, x, 1.75, w, 1.55, val, label)
            x += w + gap
        self._text(s, 0.7, 3.6, 6.0, 0.4,
                   t(self.lang, "Logique de dimensionnement", "Sizing logic"),
                   size=15, bold=True, font=self.HFONT, color=self.GREEN)
        cmap = {
            "roof": t(self.lang, "surface de toiture", "roof area"),
            "grid": t(self.lang, "puissance de raccordement", "grid connection"),
            "daytime consumption": t(self.lang, "consommation de journée", "daytime consumption"),
        }
        roof_s = f"{fmt_int(sz.cap_by_roof_kwc)} kWc" if sz.cap_by_roof_kwc > 0 else "n/a"
        grid_s = f"{fmt_int(sz.cap_by_grid_kwc)} kWc" if sz.cap_by_grid_kwc > 0 else "n/a"
        cons_s = f"{fmt_int(sz.cap_by_consumption_kwc)} kWc"
        logic = t(self.lang,
                  f"Puissance = minimum entre la toiture ({roof_s}), le raccordement ({grid_s}) "
                  f"et la consommation de journée ({cons_s}).\nFacteur limitant : "
                  f"{cmap[sz.binding_constraint]}.",
                  f"Capacity = minimum of roof ({roof_s}), grid connection ({grid_s}) and daytime "
                  f"consumption ({cons_s}).\nBinding constraint: {cmap[sz.binding_constraint]}.")
        self._text(s, 0.7, 4.1, 6.0, 1.7, logic, size=13.5, color=self.MUTED,
                   line_spacing=1.25)
        # technical card (mint)
        self._card(s, 7.0, 3.55, 5.6, 2.85)
        self._picture(s, 7.62, 3.95, 0.5, "image6.png")
        self._text(s, 8.0, 3.75, 4.4, 0.4,
                   t(self.lang, "Caractéristiques techniques", "Technical characteristics"),
                   size=15, bold=True, font=self.HFONT, color=self.GREEN)
        techs = [
            t(self.lang, "Modules monocristallins haut rendement + onduleurs string",
              "High-efficiency monocrystalline modules + string inverters"),
            t(self.lang, "Autoconsommation sur site, en complément du réseau",
              "On-site self-consumption, complementing the grid"),
            t(self.lang, "Comptage certifié de l'énergie livrée",
              "Certified metering of delivered energy"),
            t(self.lang, "Monitoring de production en temps réel",
              "Real-time production monitoring"),
        ]
        if sz.yield_source != "default":
            techs.insert(0, t(self.lang,
                              f"Productible spécifique du site : {fmt_int(sz.specific_yield)} kWh/kWc (PVGIS)",
                              f"Site-specific yield: {fmt_int(sz.specific_yield)} kWh/kWc (PVGIS)"))
        if sz.roof_area_m2:
            techs.insert(0, t(self.lang,
                              f"Surface de toiture exploitee : {fmt_int(sz.roof_area_m2)} m2",
                              f"Roof area used: {fmt_int(sz.roof_area_m2)} m2"))
        yy = 4.28
        for item in techs[:5]:
            self._disc(s, 7.4, yy + 0.05, 0.14, self.NUM)
            self._text(s, 7.72, yy - 0.03, 4.6, 0.42, item, size=12,
                       color=self.GREEN, line_spacing=1.0)
            yy += 0.40
        self._footer(s)

    def technical_feasibility(self):
        r = self.data.site_report
        s = self._slide()
        self._heading(s, t(self.lang, "Étude de faisabilité technique",
                           "Technical feasibility study"))
        verdict = r.feasibility_verdict or t(
            self.lang, "Projet techniquement réalisable sur le site.",
            "Project technically feasible on site.")
        self._card(s, 0.7, 1.55, 11.9, 0.95)
        if not self._picture(s, 1.22, 2.02, 0.6, "image11.png"):
            self._disc(s, 0.95, 1.75, 0.55, self.NUM, "OK", txt_color=self.WHITE, size=13)
        self._text(s, 1.75, 1.6, 10.7, 0.85, verdict, size=15, bold=True,
                   font=self.HFONT, color=self.GREEN, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=1.05)
        facts = []
        if r.transformer_kva:
            facts.append((f"{fmt_int(r.transformer_kva)} kVA",
                          t(self.lang, "Transformateur CIE", "CIE transformer")))
        if r.diesel_kw:
            facts.append((f"{fmt_int(r.diesel_kw)} kW",
                          t(self.lang, "Groupe de secours", "Backup generator")))
        if r.n_modules:
            facts.append((fmt_int(r.n_modules),
                          t(self.lang, "Modules PV", "PV modules")))
        if facts:
            w = 11.9 / len(facts) - 0.18
            x = 0.7
            for val, label in facts:
                self._stat_tile(s, x, 2.68, w, 1.2, val, label)
                x += w + 0.27
        bldgs = [b for b in r.buildings if b.name]
        if bldgs:
            self._text(s, 0.7, 4.12, 11.9, 0.4,
                       t(self.lang, "Répartition de la puissance par bâtiment",
                         "Power allocation per building"),
                       size=14, bold=True, color=self.GREEN)
            rows = len(bldgs) + 2
            gframe = s.shapes.add_table(rows, 3, Inches(0.7), Inches(4.5),
                                        Inches(8.6), Inches(0.31 * rows))
            table = gframe.table
            table.columns[0].width = Inches(4.0)
            table.columns[1].width = Inches(2.3)
            table.columns[2].width = Inches(2.3)
            heads = [t(self.lang, "Bâtiment", "Building"),
                     t(self.lang, "Surface", "Area"),
                     t(self.lang, "Puissance", "Capacity")]
            for j, h in enumerate(heads):
                self._set_cell(table.cell(0, j), h, size=12, bold=True,
                               color=self.WHITE, fill=self.GREEN,
                               align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
            for i, b in enumerate(bldgs, start=1):
                rfill = self.MINT if i % 2 else self.WHITE
                area = f"{fmt_int(b.area_m2)} m2" if b.area_m2 else "-"
                kwc = f"{fmt_int(b.capacity_kwc)} kWc" if b.capacity_kwc else "-"
                self._set_cell(table.cell(i, 0), b.name, size=11.5, fill=rfill)
                self._set_cell(table.cell(i, 1), area, size=11.5, fill=rfill,
                               align=PP_ALIGN.RIGHT, color=self.MUTED)
                self._set_cell(table.cell(i, 2), kwc, size=11.5, fill=rfill,
                               color=self.NUM, bold=True, align=PP_ALIGN.RIGHT)
            tot_area = sum(b.area_m2 or 0 for b in bldgs)
            tot_kwc = r.recommended_kwc or sum(b.capacity_kwc or 0 for b in bldgs)
            ri = len(bldgs) + 1
            self._set_cell(table.cell(ri, 0), "TOTAL", size=12, bold=True, fill=self.BORDER)
            self._set_cell(table.cell(ri, 1), f"{fmt_int(tot_area)} m2" if tot_area else "-",
                           size=12, bold=True, fill=self.BORDER, align=PP_ALIGN.RIGHT)
            self._set_cell(table.cell(ri, 2), f"{fmt_int(tot_kwc)} kWc", size=12,
                           bold=True, fill=self.BORDER, color=self.NUM, align=PP_ALIGN.RIGHT)
        # validation card
        self._card(s, 9.5, 4.5, 3.1, 1.95, fill=self.GREEN)
        self._text(s, 9.75, 4.68, 2.6, 0.4,
                   t(self.lang, "Audit réalisé", "Audit completed"),
                   size=13, bold=True, color=self.GOLD)
        self._text(s, 9.75, 5.14, 2.6, 0.9,
                   t(self.lang, "Visite technique et expertise du site effectuées.",
                     "Technical visit and site expertise completed."),
                   size=12, color="D9E5DD", line_spacing=1.1)
        if r.engineer_name:
            self._text(s, 9.75, 6.0, 2.6, 0.35,
                       t(self.lang, f"Valide par {r.engineer_name}",
                         f"Validated by {r.engineer_name}"),
                       size=11, italic=True, color=self.WHITE)
        self._footer(s)

    def ppa_model(self):
        s = self._slide()
        d = self.data
        copy = self.b["copy"]["ppa_model"]
        self._heading(s, t(self.lang, copy["title_fr"], copy["title_en"]))
        intro = t(self.lang,
                  f"ASPAN ENERGY finance, installe, possède et exploite la centrale sur votre site. "
                  f"{d.client.name} ne paie que l'énergie solaire consommée.",
                  f"ASPAN ENERGY finances, installs, owns and operates the plant on your site. "
                  f"{d.client.name} only pays for the solar energy consumed.")
        self._text(s, 0.7, 1.65, 11.9, 0.7, intro, size=15, color=self.MUTED,
                   line_spacing=1.2)
        pillars = copy["pillars"]
        x, y, w, h, gap = 0.7, 2.6, 2.85, 2.95, 0.18
        for p in pillars:
            self._card(s, x, y, w, h)
            self._icon_badge(s, x + 0.3, y + 0.3, 0.8, p.get("img", ""),
                             fallback=p.get("icon", ""))
            title = t(self.lang, p["title_fr"], p["title_en"])
            body = t(self.lang, p["body_fr"], p["body_en"]).format(
                years=d.contract_years, discount=d.discount_pct)
            self._text(s, x + 0.25, y + 1.2, w - 0.5, 0.7, title, size=15,
                       bold=True, font=self.HFONT, color=self.GREEN)
            self._text(s, x + 0.25, y + 1.85, w - 0.5, 0.95, body, size=12,
                       color=self.MUTED, line_spacing=1.12)
            x += w + gap
        self._footer(s)

    def _scenario_slide(self, sc, subtitle):
        s = self._slide()
        d = self.data
        self._heading(s, sc.name_fr if self.lang == "fr" else sc.name_en)
        self._text(s, 0.7, 1.6, 11.9, 0.7, subtitle, size=15, color=self.MUTED,
                   line_spacing=1.15)
        tiles = [
            (fmt_money(sc.monthly_saving, self.lang),
             t(self.lang, "Économie mensuelle", "Monthly saving"), False),
            (fmt_money(sc.annual_saving, self.lang),
             t(self.lang, "Économie annuelle estimée", "Estimated annual saving"), False),
            (fmt_money(sc.cumulative_15y_saving, self.lang),
             t(self.lang, f"Économie cumulée {d.contract_years} ans", f"Cumulative {d.contract_years}-yr saving"), True),
        ]
        x, w, gap = 0.7, 3.83, 0.2
        for val, label, emph in tiles:
            self._stat_tile(s, x, 2.5, w, 1.7, val, label, emphasis=emph)
            x += w + gap
        self._text(s, 0.7, 4.6, 11.9, 0.4,
                   t(self.lang, "Comparaison du coût du kWh", "Cost of kWh comparison"),
                   size=14, bold=True, color=self.GREEN)
        rows = [
            (t(self.lang, "Coût actuel (référence)", "Current cost (baseline)"),
             sc.baseline_cost_per_kwh, self.WARN),
            (t(self.lang, "Tarif solaire ASPAN", "ASPAN solar tariff"),
             sc.solar_tariff_per_kwh, self.NUM),
        ]
        maxv = max(r[1] for r in rows)
        yy = 5.1
        for label, val, color in rows:
            self._text(s, 0.7, yy, 3.3, 0.4, label, size=12.5, color=self.GREEN,
                       anchor=MSO_ANCHOR.MIDDLE)
            bw = max(7.0 * val / maxv, 0.8)
            self._card(s, 4.1, yy + 0.03, bw, 0.34, fill=color, radius=0.4)
            self._text(s, 4.1 + bw + 0.15, yy, 1.6, 0.4, f"{val:.2f}", size=12.5,
                       bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
            yy += 0.55
        self._text(s, 0.7, 6.4, 11.9, 0.45,
                   t(self.lang, f"Soit -{sc.saving_pct:.0f}% sur le coût du kWh substitué par le solaire.",
                     f"That is -{sc.saving_pct:.0f}% on the cost of each kWh replaced by solar."),
                   size=13.5, color=self.NUM, bold=True)
        self._footer(s)

    def scenario_1(self):
        d = self.data
        sub = t(self.lang,
                f"La consommation de journée est aujourd'hui fournie par le réseau CIE "
                f"({d.cie_day_tariff:.2f} FCFA/kWh). Le solaire la remplace au tarif ASPAN "
                f"({d.solar_tariff:.2f} FCFA/kWh).",
                f"Daytime consumption is currently supplied by the CIE grid "
                f"({d.cie_day_tariff:.2f} FCFA/kWh). Solar replaces it at the ASPAN tariff "
                f"({d.solar_tariff:.2f} FCFA/kWh).")
        self._scenario_slide(d.scenarios[0], sub)

    def scenario_2(self):
        d = self.data
        a = assumptions()
        diesel = a["diesel"]["fuel_price_fcfa_per_litre"] * a["diesel"]["litres_per_kwh"]
        sub = t(self.lang,
                f"Hypothèse : 50% réseau / 50% groupe électrogène (gasoil {a['diesel']['fuel_price_fcfa_per_litre']:.0f} "
                f"FCFA/l, {a['diesel']['litres_per_kwh']} l/kWh, soit {diesel:.0f} FCFA/kWh). "
                f"Le solaire remplace ce mix à {d.solar_tariff:.2f} FCFA/kWh.",
                f"Assumption: 50% grid / 50% diesel generator (diesel {a['diesel']['fuel_price_fcfa_per_litre']:.0f} "
                f"FCFA/l, {a['diesel']['litres_per_kwh']} l/kWh, i.e. {diesel:.0f} FCFA/kWh). "
                f"Solar replaces this mix at {d.solar_tariff:.2f} FCFA/kWh.")
        self._scenario_slide(d.scenarios[1], sub)

    def savings_chart(self):
        s = self._slide()
        d = self.data
        self._heading(s, t(self.lang, f"Économies cumulées sur {d.contract_years} ans",
                           f"Cumulative savings over {d.contract_years} years"))
        years = list(range(1, d.contract_years + 1))
        chart_data = CategoryChartData()
        chart_data.categories = [str(y) for y in years]
        for sc in d.scenarios:
            name = sc.name_fr if self.lang == "fr" else sc.name_en
            chart_data.add_series(
                name, [round(sc.annual_saving * y / 1_000_000, 1) for y in years])
        gframe = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.7),
                                    Inches(1.8), Inches(12.0), Inches(4.7), chart_data)
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.color.rgb = _rgb(self.GREEN)
        chart.legend.font.size = Pt(12)
        series_colors = [self.NUM, self.GOLD]
        for i, ser in enumerate(chart.series):
            ser.format.line.color.rgb = _rgb(series_colors[i % 2])
            ser.format.line.width = Pt(2.5)
        for axis in (chart.category_axis, chart.value_axis):
            axis.tick_labels.font.color.rgb = _rgb(self.MUTED)
            axis.tick_labels.font.size = Pt(10)
        self._text(s, 0.7, 6.55, 11.9, 0.4,
                   t(self.lang, "Économies cumulées en millions de FCFA (HT).",
                     "Cumulative savings in millions of FCFA (excl. tax)."),
                   size=11, color=self.MUTED)
        self._footer(s)

    def financial_analysis(self):
        s = self._slide()
        d = self.data
        f = d.financial
        self._heading(s, t(self.lang, "Analyse financière", "Financial analysis"))
        tiles = [
            (fmt_money(f.npv_savings_escalated, self.lang),
             t(self.lang, f"VAN des économies ({d.contract_years} ans)",
               f"NPV of savings ({d.contract_years} yr)"), True),
            (f"{f.solar_coverage_pct:.0f}%",
             t(self.lang, "de la consommation couverte", "of consumption covered"), False),
            (fmt_money(f.annual_saving, self.lang),
             t(self.lang, "Économie annuelle (an 1)", "Annual saving (yr 1)"), False),
        ]
        x, w, gap = 0.7, 3.83, 0.2
        for val, label, emph in tiles:
            self._stat_tile(s, x, 1.7, w, 1.5, val, label, emphasis=emph)
            x += w + gap
        self._text(s, 0.7, 3.55, 5.4, 0.4,
                   t(self.lang, "Hypothèses & impact", "Assumptions & impact"),
                   size=15, bold=True, font=self.HFONT, color=self.GREEN)
        body = t(self.lang,
                 f"Taux d'actualisation : {f.discount_rate*100:.0f}%  -  escalade tarif CIE : "
                 f"{f.tariff_escalation*100:.0f}%/an.\n"
                 f"Coût actuel de cette énergie : {fmt_money(f.current_annual_cost, self.lang)}/an.\n"
                 f"Coût sous PPA ASPAN : {fmt_money(f.solar_annual_cost, self.lang)}/an.\n"
                 f"Impact CO2 évité : {fmt_int(f.co2_tonnes_per_year)} t/an "
                 f"(~{fmt_int(f.co2_tonnes_horizon)} t sur {d.contract_years} ans).",
                 f"Discount rate: {f.discount_rate*100:.0f}%  -  CIE tariff escalation: "
                 f"{f.tariff_escalation*100:.0f}%/yr.\n"
                 f"Current cost of this energy: {fmt_money(f.current_annual_cost, self.lang)}/yr.\n"
                 f"Cost under the ASPAN PPA: {fmt_money(f.solar_annual_cost, self.lang)}/yr.\n"
                 f"CO2 avoided: {fmt_int(f.co2_tonnes_per_year)} t/yr "
                 f"(~{fmt_int(f.co2_tonnes_horizon)} t over {d.contract_years} yr).")
        self._text(s, 0.7, 4.05, 5.4, 2.3, body, size=13, color=self.MUTED,
                   line_spacing=1.3)
        years = [cy for cy in f.cashflow if cy.year in (1, 3, 5, 7, 9, 11, 13, 15)]
        rows = len(years) + 1
        gframe = s.shapes.add_table(rows, 3, Inches(6.4), Inches(3.62),
                                    Inches(6.2), Inches(2.95))
        table = gframe.table
        table.columns[0].width = Inches(1.4)
        table.columns[1].width = Inches(2.4)
        table.columns[2].width = Inches(2.4)
        headers = [t(self.lang, "Année", "Year"),
                   t(self.lang, "Économie/an", "Saving/yr"),
                   t(self.lang, "Cumul", "Cumulative")]
        for j, htxt in enumerate(headers):
            self._set_cell(table.cell(0, j), htxt, size=12, bold=True,
                           color=self.WHITE, fill=self.GREEN,
                           align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
        for i, cy in enumerate(years, start=1):
            rfill = self.MINT if i % 2 else self.WHITE
            self._set_cell(table.cell(i, 0), str(cy.year), size=12, fill=rfill,
                           color=self.GREEN)
            self._set_cell(table.cell(i, 1), fmt_money(cy.annual_saving, self.lang),
                           size=12, fill=rfill, color=self.MUTED, align=PP_ALIGN.RIGHT)
            self._set_cell(table.cell(i, 2), fmt_money(cy.cumulative_saving, self.lang),
                           size=12, fill=rfill, color=self.NUM, bold=True,
                           align=PP_ALIGN.RIGHT)
        self._footer(s)

    def services(self):
        s = self._slide()
        copy = self.b["copy"]["services"]
        d = self.data
        title = t(self.lang, copy["title_fr"], copy["title_en"]).format(years=d.contract_years)
        self._heading(s, title)
        items = copy["items"]
        x0, y0, w, h, gx, gyy = 0.7, 1.85, 3.86, 1.95, 0.2, 0.22
        for idx, it in enumerate(items):
            col, row = idx % 3, idx // 3
            x = x0 + col * (w + gx)
            y = y0 + row * (h + gyy)
            self._card(s, x, y, w, h)
            self._icon_badge(s, x + 0.3, y + 0.3, 0.6, it.get("img", ""),
                             fallback=str(idx + 1))
            self._text(s, x + 1.05, y + 0.32, w - 1.25, 0.6,
                       t(self.lang, it["fr"], it["en"]), size=14.5, bold=True,
                       font=self.HFONT, color=self.GREEN)
            self._text(s, x + 0.3, y + 1.05, w - 0.6, 0.8,
                       t(self.lang, it["desc_fr"], it["desc_en"]), size=12,
                       color=self.MUTED, line_spacing=1.12)
        self._footer(s)

    def next_steps(self):
        s = self._slide(dark=True, decor="closing")
        d = self.data
        c = self.b["company"]
        copy = self.b["copy"]["next_steps"]
        self._heading(s, t(self.lang, copy["title_fr"], copy["title_en"]), dark=True)
        steps = copy["steps"]
        done_count = 2 if d.site_report else 0
        y = 1.85
        for i, st in enumerate(steps):
            is_done = i < done_count
            self._disc(s, 0.7, y, 0.6, self.GOLD, str(i + 1),
                       txt_color=self.GREEN, size=18)
            label = t(self.lang, st["fr"], st["en"]).format(years=d.contract_years)
            if is_done:
                label += t(self.lang, "  (effectue)", "  (done)")
            self._text(s, 1.5, y + 0.02, 10.5, 0.6, label, size=15,
                       color=self.WHITE, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.82
        # contact block bottom-left
        self._text(s, 0.7, 5.55, 7.0, 0.4, c["contact_name"], size=17, bold=True,
                   color=self.WHITE)
        self._text(s, 0.7, 6.0, 8.0, 0.35,
                   t(self.lang, c["contact_role_fr"], c["contact_role_en"]),
                   size=13, color="D9E5DD")
        self._text(s, 0.7, 6.4, 9.0, 0.35,
                   f"{c['contact_email']}   -   {c['contact_address']}",
                   size=13, color="D9E5DD")
        self._text(s, 0.7, 6.78, 9.0, 0.35,
                   t(self.lang,
                     f"Offre indicative valable {c['offer_validity_days']} jours - Ref. {d.reference}",
                     f"Indicative offer valid for {c['offer_validity_days']} days - Ref. {d.reference}"),
                   size=11, italic=True, color="A9C2B4")

    def build(self) -> Presentation:
        # Slide order mirrors Aspan's own deck:
        # cover -> PPA model -> project/solution -> scenarios -> services -> next steps.
        # Our extra (assignment-required / value-add) slides are slotted in around them.
        self.cover()
        self.executive_summary()
        self.current_situation()
        self.ppa_model()
        self.solar_solution()
        if self.data.site_report:
            self.technical_feasibility()
        self.scenario_1()
        self.scenario_2()
        self.savings_chart()
        self.financial_analysis()
        self.services()
        self.next_steps()
        return self.prs


def build_pptx(data: ProposalData, output_path: str) -> str:
    """Build the deck and save it. Returns the output path."""
    deck = DeckBuilder(data)
    prs = deck.build()
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
